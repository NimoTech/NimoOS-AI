import pytest


def test_unsupported_extension_returns_error(tmp_path):
    from attachments.extract import extract_to_markdown
    p = tmp_path / "x.xyz"
    p.write_bytes(b"hello")
    result = extract_to_markdown(str(p), "xyz",
                                 max_chars=1000,
                                 max_uncompressed_bytes=10_000_000)
    assert result == {"ok": False, "error": "unsupported"}


def test_missing_library_returns_not_installed(tmp_path, monkeypatch):
    """If the per-format library can't be imported, the dispatcher reports
    not_installed for that extension only."""
    from attachments import extract

    # Force the pdf branch's import to fail by stubbing builtins.__import__.
    real_import = __builtins__["__import__"] if isinstance(__builtins__, dict) else __builtins__.__import__

    def fake_import(name, *args, **kwargs):
        if name == "pypdf":
            raise ImportError("test forcibly hides pypdf")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr("builtins.__import__", fake_import)

    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4\n")
    result = extract.extract_to_markdown(str(p), "pdf",
                                         max_chars=1000,
                                         max_uncompressed_bytes=10_000_000)
    assert result == {"ok": False, "error": "not_installed"}


import os


def _fixture(name: str) -> str:
    return os.path.join(os.path.dirname(__file__), "fixtures", name)


def test_pdf_happy_path():
    from attachments.extract import extract_to_markdown
    result = extract_to_markdown(_fixture("tiny.pdf"), "pdf",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert "Hello PDF" in result["markdown"]
    assert result["pages"] == 1
    assert result["extractor"] == "pypdf"
    assert result["truncated"] is False


def test_pdf_truncation_sets_flag(monkeypatch):
    """With max_chars below the extracted length, extractor stops early."""
    from attachments.extract import extract_to_markdown
    # max_chars=5 forces truncation since "Hello PDF" is 9 chars.
    result = extract_to_markdown(_fixture("tiny.pdf"), "pdf",
                                 max_chars=5,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert result["truncated"] is True
    assert len(result["markdown"]) <= 5


def test_pdf_garbage_returns_parse_error(tmp_path):
    from attachments.extract import extract_to_markdown
    p = tmp_path / "junk.pdf"
    p.write_bytes(b"this is not a pdf at all")
    result = extract_to_markdown(str(p), "pdf",
                                 max_chars=1000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is False
    assert result["error"] == "parse_error"


def test_pdf_encrypted_returns_encrypted(tmp_path, monkeypatch):
    """We don't need an actual encrypted PDF — monkeypatch pypdf to claim
    is_encrypted=True and refuse decrypt('')."""
    from attachments import extract
    import pypdf

    class FakeReader:
        is_encrypted = True
        pages: list = []
        def __init__(self, *_a, **_kw): pass
        def decrypt(self, _pw): return 0  # 0 = failure in pypdf API

    monkeypatch.setattr(pypdf, "PdfReader", FakeReader)
    p = tmp_path / "enc.pdf"
    p.write_bytes(b"%PDF-1.4\n")  # presence only; FakeReader ignores content
    result = extract.extract_to_markdown(str(p), "pdf",
                                         max_chars=1000,
                                         max_uncompressed_bytes=10_000_000)
    assert result == {"ok": False, "error": "encrypted"}


def test_pdf_empty_extraction_returns_empty_scanned(tmp_path, monkeypatch):
    """Simulate a scanned PDF: pypdf yields whitespace-only text."""
    from attachments import extract
    import pypdf

    class FakePage:
        def extract_text(self): return "   \n  "

    class FakeReader:
        is_encrypted = False
        pages = [FakePage()]
        def __init__(self, *_a, **_kw): pass

    monkeypatch.setattr(pypdf, "PdfReader", FakeReader)
    p = tmp_path / "scan.pdf"
    p.write_bytes(b"%PDF-1.4\n")
    result = extract.extract_to_markdown(str(p), "pdf",
                                         max_chars=1000,
                                         max_uncompressed_bytes=10_000_000)
    assert result == {"ok": False, "error": "empty_scanned"}


def test_pdf_unexpected_exception_caught(tmp_path, monkeypatch):
    """Even an exception type we never anticipated must not escape."""
    from attachments import extract
    import pypdf

    class ExplodingReader:
        def __init__(self, *_a, **_kw):
            raise RecursionError("simulated pypdf pathology")

    monkeypatch.setattr(pypdf, "PdfReader", ExplodingReader)
    p = tmp_path / "boom.pdf"
    p.write_bytes(b"%PDF-1.4\n")
    result = extract.extract_to_markdown(str(p), "pdf",
                                         max_chars=1000,
                                         max_uncompressed_bytes=10_000_000)
    assert result == {"ok": False, "error": "parse_error"}


def _build_docx(tmp_path, paragraphs):
    """Build an in-memory .docx using python-docx and return its path."""
    import docx
    d = docx.Document()
    for p in paragraphs:
        d.add_paragraph(p)
    p = tmp_path / "doc.docx"
    d.save(str(p))
    return str(p)


def test_docx_happy_path(tmp_path):
    from attachments.extract import extract_to_markdown
    path = _build_docx(tmp_path, ["First paragraph.", "Second paragraph."])
    result = extract_to_markdown(path, "docx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert "First paragraph." in result["markdown"]
    assert "Second paragraph." in result["markdown"]
    assert result["extractor"] == "python-docx"
    assert result["truncated"] is False


def test_docx_truncation(tmp_path):
    from attachments.extract import extract_to_markdown
    big = "x" * 500
    path = _build_docx(tmp_path, [big, big, big])
    result = extract_to_markdown(path, "docx",
                                 max_chars=100,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert result["truncated"] is True
    assert len(result["markdown"]) <= 100


def test_docx_garbage_returns_parse_error(tmp_path):
    from attachments.extract import extract_to_markdown
    p = tmp_path / "junk.docx"
    p.write_bytes(b"not a zip at all")
    result = extract_to_markdown(str(p), "docx",
                                 max_chars=1000,
                                 max_uncompressed_bytes=10_000_000)
    # Garbage that isn't even a ZIP fails the zip-bomb precheck (BadZipFile
    # → returns False → zip_bomb branch). If your implementation chose a
    # different signaling for "not even a zip", adjust this assertion to
    # match — but document the choice in the helper's docstring.
    assert result == {"ok": False, "error": "zip_bomb"} or result == {"ok": False, "error": "parse_error"}


def test_docx_zip_bomb_precheck(tmp_path):
    """A ZIP whose summed uncompressed size exceeds the cap is rejected
    before python-docx ever touches it."""
    import zipfile
    from attachments.extract import extract_to_markdown
    p = tmp_path / "bomb.docx"
    # Two entries totalling ~3000 bytes uncompressed; cap of 1000 rejects.
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("a.xml", "A" * 2000)
        z.writestr("b.xml", "B" * 1000)
    result = extract_to_markdown(str(p), "docx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=1000)
    assert result == {"ok": False, "error": "zip_bomb"}


def _build_xlsx(tmp_path, sheets):
    """sheets: list of (title, list-of-rows). Returns path."""
    import openpyxl
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for title, rows in sheets:
        ws = wb.create_sheet(title=title)
        for row in rows:
            ws.append(row)
    p = tmp_path / "book.xlsx"
    wb.save(str(p))
    return str(p)


def test_xlsx_happy_path(tmp_path):
    from attachments.extract import extract_to_markdown
    path = _build_xlsx(tmp_path, [
        ("Sheet1", [["Name", "Age"], ["Alice", 30], ["Bob", 25]]),
    ])
    result = extract_to_markdown(path, "xlsx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert "## Sheet1" in result["markdown"]
    assert "Alice" in result["markdown"]
    assert "Bob" in result["markdown"]
    assert result["pages"] == 1  # sheet count
    assert result["extractor"] == "openpyxl"


def test_xlsx_multiple_sheets(tmp_path):
    from attachments.extract import extract_to_markdown
    path = _build_xlsx(tmp_path, [
        ("Numbers", [[1, 2], [3, 4]]),
        ("Letters", [["a", "b"], ["c", "d"]]),
    ])
    result = extract_to_markdown(path, "xlsx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert "## Numbers" in result["markdown"]
    assert "## Letters" in result["markdown"]
    assert result["pages"] == 2


def test_xlsx_data_only_falls_back_to_formula(tmp_path):
    """A workbook whose cells are all formulas with no cached calc values
    (the canonical Python-generated case) triggers the data_only=False
    fallback so the model gets formula text instead of all-None rows."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Calc"
    # Formulas only — openpyxl never writes a cached <v>, so data_only=True
    # returns None for every cell.
    ws["A1"] = "=1+1"
    ws["A2"] = "=2+2"
    p = tmp_path / "formula.xlsx"
    wb.save(str(p))

    from attachments.extract import extract_to_markdown
    result = extract_to_markdown(str(p), "xlsx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    # Fallback path emits the formula text rather than blank cells.
    assert "=1+1" in result["markdown"]
    assert "=2+2" in result["markdown"]


def test_xlsx_zip_bomb_precheck(tmp_path):
    import zipfile
    from attachments.extract import extract_to_markdown
    p = tmp_path / "bomb.xlsx"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("a.xml", "A" * 5000)
    result = extract_to_markdown(str(p), "xlsx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=1000)
    assert result == {"ok": False, "error": "zip_bomb"}


def test_xlsx_garbage_returns_parse_error(tmp_path):
    from attachments.extract import extract_to_markdown
    p = tmp_path / "junk.xlsx"
    p.write_bytes(b"not a zip")
    result = extract_to_markdown(str(p), "xlsx",
                                 max_chars=1000,
                                 max_uncompressed_bytes=10_000_000)
    # Non-ZIP fails _zipbomb_check → zip_bomb (consistent with docx behavior)
    assert result == {"ok": False, "error": "zip_bomb"} or result == {"ok": False, "error": "parse_error"}


def _build_pptx(tmp_path, slides_text):
    """slides_text: list[list[str]] — per slide, list of paragraphs."""
    from pptx import Presentation
    p = Presentation()
    layout = p.slide_layouts[5]  # title only
    for paragraphs in slides_text:
        slide = p.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = paragraphs[0] if paragraphs else ""
        if len(paragraphs) > 1:
            txbox = slide.shapes.add_textbox(left=0, top=0, width=100, height=100)
            tf = txbox.text_frame
            tf.text = paragraphs[1]
            for extra in paragraphs[2:]:
                tf.add_paragraph().text = extra
    out = tmp_path / "deck.pptx"
    p.save(str(out))
    return str(out)


def test_pptx_happy_path(tmp_path):
    from attachments.extract import extract_to_markdown
    path = _build_pptx(tmp_path, [
        ["Title One", "Body of slide 1"],
        ["Title Two", "Body of slide 2", "Another paragraph"],
    ])
    result = extract_to_markdown(path, "pptx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert "Title One" in result["markdown"]
    assert "Title Two" in result["markdown"]
    assert "Body of slide 1" in result["markdown"]
    assert result["pages"] == 2
    assert result["extractor"] == "python-pptx"


def test_pptx_truncation(tmp_path):
    from attachments.extract import extract_to_markdown
    path = _build_pptx(tmp_path, [["T", "x" * 5000]] * 3)
    result = extract_to_markdown(path, "pptx",
                                 max_chars=200,
                                 max_uncompressed_bytes=10_000_000)
    assert result["ok"] is True
    assert result["truncated"] is True
    assert len(result["markdown"]) <= 200


def test_pptx_zip_bomb_precheck(tmp_path):
    import zipfile
    from attachments.extract import extract_to_markdown
    p = tmp_path / "bomb.pptx"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("a.xml", "A" * 5000)
    result = extract_to_markdown(str(p), "pptx",
                                 max_chars=10_000,
                                 max_uncompressed_bytes=1000)
    assert result == {"ok": False, "error": "zip_bomb"}


def test_pptx_garbage_returns_parse_error(tmp_path):
    from attachments.extract import extract_to_markdown
    p = tmp_path / "junk.pptx"
    p.write_bytes(b"not a zip")
    result = extract_to_markdown(str(p), "pptx",
                                 max_chars=1000,
                                 max_uncompressed_bytes=10_000_000)
    # Non-ZIP fails _zipbomb_check → zip_bomb (consistent with docx/xlsx)
    assert result == {"ok": False, "error": "zip_bomb"} or result == {"ok": False, "error": "parse_error"}
